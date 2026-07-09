#!/usr/bin/env python3
"""
classroom_analyzer_helper.py
SLC 課例分析輔助腳本
用法:
  uv run scripts/classroom_analyzer_helper.py generate-slides --analysis output/analysis.txt --output output/slides.pptx --style learning
  uv run scripts/classroom_analyzer_helper.py generate-image   --text output/concept.txt   --output output/concept_post.png --style learning
"""

import argparse
import sys
import os
import re
import textwrap
from pathlib import Path

# ── 色彩主題 ─────────────────────────────────────────────────────────────────
THEMES = {
    "learning": {
        "bg":    (248, 246, 240),   # 有機茶白
        "title": (42,  82,  52),    # 深松針綠
        "body":  (72,  80,  72),    # 溫潤石板灰
        "accent":(107, 142, 107),   # 淡松針綠
        "line":  (180, 200, 180),
    },
    "pastel": {
        "bg":    (255, 253, 248),
        "title": (90,  55,  30),
        "body":  (40,  40,  40),
        "accent":(200, 150, 100),
        "line":  (210, 185, 160),
    },
    "blue": {
        "bg":    (235, 242, 250),
        "title": (20,  50,  100),
        "body":  (50,  70,  110),
        "accent":(70,  120, 200),
        "line":  (150, 180, 220),
    },
    "modern": {
        "bg":    (245, 245, 245),
        "title": (20,  20,  20),
        "body":  (80,  80,  80),
        "accent":(120, 120, 120),
        "line":  (180, 180, 180),
    },
}

# ── 解析分析文件 ──────────────────────────────────────────────────────────────
def parse_analysis(text: str) -> list[dict]:
    """
    解析 analysis.txt。支援兩種格式：
    格式一：## 第N頁 標題\n描述/詮釋/反思 section
    格式二：純段落（每個段落自動變成一張投影片）
    """
    slides = []

    # 嘗試結構化格式
    slide_pattern = re.compile(
        r"(?:#{1,3}\s*(?:第\s*\d+\s*[頁張]|Slide\s*\d+)[：:、\s]*)(.+?)$",
        re.MULTILINE
    )
    matches = list(slide_pattern.finditer(text))

    if matches:
        for i, m in enumerate(matches):
            title = m.group(1).strip()
            start = m.end()
            end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
            body = text[start:end].strip()
            slides.append({"title": title, "body": body})
    else:
        # 段落模式
        paragraphs = [p.strip() for p in re.split(r"\n{2,}", text) if p.strip()]
        for i, p in enumerate(paragraphs):
            lines = p.splitlines()
            title = lines[0].lstrip("#").strip() if lines else f"第 {i+1} 頁"
            body = "\n".join(lines[1:]).strip() if len(lines) > 1 else ""
            slides.append({"title": title, "body": body})

    return slides or [{"title": "課例分析", "body": text.strip()}]


# ── 產生 PowerPoint ───────────────────────────────────────────────────────────
def generate_pptx(slides: list[dict], output_path: str, style: str, images_dir: str | None = None):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        print("[ERROR] python-pptx required: uv pip install python-pptx")
        sys.exit(1)

    theme = THEMES.get(style, THEMES["learning"])
    bg_rgb    = RGBColor(*theme["bg"])
    title_rgb = RGBColor(*theme["title"])
    body_rgb  = RGBColor(*theme["body"])
    accent_rgb= RGBColor(*theme["accent"])

    W, H = Inches(13.333), Inches(7.5)  # 16:9

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    blank_layout = prs.slide_layouts[6]  # 空白版型

    def add_slide(title_text: str, body_text: str, img_path: str | None = None):
        slide = prs.slides.add_slide(blank_layout)

        # 背景色
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = bg_rgb

        # 裝飾橫線
        from pptx.util import Pt as UPt
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.5), W - Inches(1), Emu(3))
        line.fill.background()
        line.line.color.rgb = RGBColor(*theme["line"])
        line.line.width = Pt(1)

        # 圖片區域設定
        has_img = img_path and os.path.exists(img_path)
        text_w = W - Inches(1) if not has_img else W - Inches(5.5)
        text_x = Inches(0.5)
        img_x  = W - Inches(4.8)

        # 標題
        title_box = slide.shapes.add_textbox(text_x, Inches(0.25), text_w, Inches(1.1))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title_text
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = title_rgb
        run.font.name = "微軟正黑體"

        # 內文
        body_box = slide.shapes.add_textbox(text_x, Inches(1.7), text_w, Inches(5.3))
        tf2 = body_box.text_frame
        tf2.word_wrap = True
        for line_text in body_text.split("\n"):
            p2 = tf2.add_paragraph()
            p2.alignment = PP_ALIGN.LEFT
            run2 = p2.add_run()
            run2.text = line_text
            run2.font.size = Pt(20)
            run2.font.color.rgb = body_rgb
            run2.font.name = "微軟正黑體"

        # 插圖（若存在）
        if has_img:
            try:
                slide.shapes.add_picture(img_path, img_x, Inches(1.3), Inches(4.5), Inches(5.8))
            except Exception as e:
                print(f"⚠️ 插圖載入失敗 ({img_path}): {e}")

    # 封面投影片
    cover = prs.slides.add_slide(blank_layout)
    bg = cover.background.fill
    bg.solid()
    bg.fore_color.rgb = bg_rgb

    cov_title = cover.shapes.add_textbox(Inches(1.5), Inches(2.5), W - Inches(3), Inches(1.5))
    tf = cov_title.text_frame
    p = tf.paragraphs[0]
    from pptx.enum.text import PP_ALIGN
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = slides[0]["title"] if slides else "課例研究分析"
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = title_rgb
    run.font.name = "微軟正黑體"

    # 內容投影片
    for i, s in enumerate(slides):
        img_path = None
        if images_dir:
            candidate = os.path.join(images_dir, f"slide_{i+1}.png")
            if os.path.exists(candidate):
                img_path = candidate
        add_slide(s["title"], s["body"], img_path)

    # 儲存（若被鎖定，自動存為 _copy）
    out = Path(output_path)
    try:
        prs.save(str(out))
        print(f"[OK] PPTX saved: {out}")
    except PermissionError:
        copy_out = out.with_stem(out.stem + "_copy")
        prs.save(str(copy_out))
        print(f"[WARN] File locked, saved as: {copy_out}")
        out = copy_out

    return str(out)


# ── 產生 HTML 簡報 ────────────────────────────────────────────────────────────
def generate_html(slides: list[dict], pptx_path: str, style: str, images_dir: str | None = None):
    import base64

    theme = THEMES.get(style, THEMES["learning"])
    r, g, b = theme["bg"]
    tr, tg, tb = theme["title"]
    br, bg2, bb = theme["body"]
    ar, ag, ab = theme["accent"]

    def rgb_to_hex(r, g, b):
        return f"#{r:02x}{g:02x}{b:02x}"

    bg_hex    = rgb_to_hex(*theme["bg"])
    title_hex = rgb_to_hex(*theme["title"])
    body_hex  = rgb_to_hex(*theme["body"])
    accent_hex= rgb_to_hex(*theme["accent"])
    line_hex  = rgb_to_hex(*theme["line"])

    def img_to_b64(path):
        try:
            with open(path, "rb") as f:
                return base64.b64encode(f.read()).decode()
        except Exception:
            return None

    slide_html = []
    for i, s in enumerate(slides):
        img_b64 = None
        if images_dir:
            candidate = os.path.join(images_dir, f"slide_{i+1}.png")
            if os.path.exists(candidate):
                img_b64 = img_to_b64(candidate)

        img_section = ""
        if img_b64:
            img_section = f'<img class="slide-img" src="data:image/png;base64,{img_b64}" alt="插圖">'

        body_html = "".join(f"<p>{line}</p>" for line in s["body"].split("\n") if line.strip())

        slide_html.append(f"""
        <div class="slide" id="slide-{i+1}">
          <div class="slide-content {'has-img' if img_b64 else ''}">
            <div class="text-area">
              <h2>{s["title"]}</h2>
              <div class="body-text">{body_html}</div>
            </div>
            {img_section}
          </div>
          <div class="slide-number">{i+1} / {len(slides)}</div>
        </div>
        """)

    slides_str = "\n".join(slide_html)

    html = f"""<!DOCTYPE html>
<html lang="zh-TW">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>課例研究分析簡報</title>
<style>
  * {{ box-sizing: border-box; margin: 0; padding: 0; }}
  body {{
    background: #1a1a2e;
    display: flex;
    flex-direction: column;
    align-items: center;
    justify-content: center;
    min-height: 100vh;
    font-family: "微軟正黑體", "Noto Sans TC", sans-serif;
  }}
  .slideshow {{
    width: min(1280px, 95vw);
    aspect-ratio: 16/9;
    position: relative;
    overflow: hidden;
    border-radius: 12px;
    box-shadow: 0 20px 60px rgba(0,0,0,0.5);
  }}
  .slide {{
    position: absolute;
    inset: 0;
    background: {bg_hex};
    display: none;
    padding: 3% 4%;
    flex-direction: column;
    justify-content: flex-start;
  }}
  .slide.active {{ display: flex; }}
  .slide-content {{
    display: flex;
    flex-direction: row;
    gap: 4%;
    height: 88%;
    align-items: flex-start;
  }}
  .text-area {{
    flex: 1;
    display: flex;
    flex-direction: column;
    gap: 1.2em;
  }}
  h2 {{
    color: {title_hex};
    font-size: clamp(1.4rem, 3.2vw, 2.4rem);
    font-weight: 700;
    border-bottom: 3px solid {accent_hex};
    padding-bottom: 0.4em;
    line-height: 1.3;
  }}
  .body-text {{ color: {body_hex}; font-size: clamp(0.9rem, 1.8vw, 1.25rem); line-height: 1.9; }}
  .body-text p {{ margin-bottom: 0.6em; }}
  .slide-img {{
    width: 38%;
    max-height: 90%;
    object-fit: contain;
    border-radius: 8px;
    box-shadow: 0 4px 20px rgba(0,0,0,0.15);
    flex-shrink: 0;
  }}
  .slide-number {{
    position: absolute;
    bottom: 1.5%;
    right: 2%;
    color: {accent_hex};
    font-size: clamp(0.7rem, 1.2vw, 0.9rem);
    opacity: 0.7;
  }}
  .controls {{
    display: flex;
    gap: 16px;
    margin-top: 20px;
    align-items: center;
  }}
  button {{
    background: {accent_hex};
    color: white;
    border: none;
    padding: 10px 28px;
    border-radius: 6px;
    font-size: 1rem;
    font-family: inherit;
    cursor: pointer;
    transition: opacity 0.2s;
  }}
  button:hover {{ opacity: 0.8; }}
  .progress {{
    color: #aaa;
    font-family: "微軟正黑體", sans-serif;
    font-size: 0.95rem;
    min-width: 80px;
    text-align: center;
  }}
</style>
</head>
<body>
<div class="slideshow" id="slideshow">
{slides_str}
</div>
<div class="controls">
  <button onclick="changeSlide(-1)">◀ 上一頁</button>
  <span class="progress" id="progress">1 / {len(slides)}</span>
  <button onclick="changeSlide(1)">下一頁 ▶</button>
</div>
<script>
  let cur = 0;
  const slides = document.querySelectorAll('.slide');
  function show(n) {{
    slides[cur].classList.remove('active');
    cur = (n + slides.length) % slides.length;
    slides[cur].classList.add('active');
    document.getElementById('progress').textContent = (cur+1) + ' / ' + slides.length;
  }}
  function changeSlide(d) {{ show(cur + d); }}
  document.addEventListener('keydown', e => {{
    if (e.key === 'ArrowRight' || e.key === ' ') changeSlide(1);
    if (e.key === 'ArrowLeft') changeSlide(-1);
  }});
  slides[0].classList.add('active');
</script>
</body>
</html>"""

    html_path = Path(pptx_path).with_suffix(".html")
    try:
        html_path.write_text(html, encoding="utf-8")
        print(f"[OK] HTML saved: {html_path}")
    except Exception as e:
        print(f"[WARN] HTML save failed: {e}")


# ── 產生社群概念圖 ────────────────────────────────────────────────────────────
def generate_image_cmd(text_path: str, output_path: str, style: str):
    try:
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        print("[ERROR] Pillow required: uv pip install Pillow")
        sys.exit(1)

    theme = THEMES.get(style, THEMES["learning"])
    W, H = 1200, 630

    img = Image.new("RGB", (W, H), color=theme["bg"])
    draw = ImageDraw.Draw(img)

    content = Path(text_path).read_text(encoding="utf-8")
    lines = content.strip().splitlines()

    # 嘗試載入字型（fallback 到 default）
    try:
        title_font = ImageFont.truetype("C:/Windows/Fonts/msjhbd.ttc", 48)
        body_font  = ImageFont.truetype("C:/Windows/Fonts/msjh.ttc", 28)
    except Exception:
        title_font = ImageFont.load_default()
        body_font  = ImageFont.load_default()

    y = 60
    if lines:
        draw.text((80, y), lines[0], fill=theme["title"], font=title_font)
        y += 80

    for line in lines[1:]:
        if y > H - 60:
            break
        wrapped = textwrap.fill(line, width=38)
        for wl in wrapped.splitlines():
            draw.text((80, y), wl, fill=theme["body"], font=body_font)
            y += 44

    # 裝飾線
    draw.rectangle([60, 130, 60, H - 60], fill=theme["accent"], width=4)

    img.save(output_path)
    print(f"[OK] Concept image saved: {output_path}")


# ── CLI 主入口 ────────────────────────────────────────────────────────────────
def main():
    parser = argparse.ArgumentParser(description="SLC 課例分析輔助工具")
    sub = parser.add_subparsers(dest="cmd")

    # generate-slides
    p_slides = sub.add_parser("generate-slides")
    p_slides.add_argument("--analysis", required=True)
    p_slides.add_argument("--output",   required=True)
    p_slides.add_argument("--style",    default="learning", choices=THEMES.keys())
    p_slides.add_argument("--images",   default=None, help="插圖目錄（如 output/images）")

    # generate-image
    p_img = sub.add_parser("generate-image")
    p_img.add_argument("--text",   required=True)
    p_img.add_argument("--output", required=True)
    p_img.add_argument("--style",  default="learning", choices=THEMES.keys())

    args = parser.parse_args()

    if args.cmd == "generate-slides":
        text = Path(args.analysis).read_text(encoding="utf-8")
        slides = parse_analysis(text)
        print(f"[INFO] Parsed {len(slides)} slides")

        # 自動偵測 images 目錄
        images_dir = args.images
        if not images_dir:
            candidate = str(Path(args.output).parent / "images")
            if os.path.isdir(candidate):
                images_dir = candidate

        pptx_path = generate_pptx(slides, args.output, args.style, images_dir)
        generate_html(slides, pptx_path, args.style, images_dir)

    elif args.cmd == "generate-image":
        generate_image_cmd(args.text, args.output, args.style)

    else:
        parser.print_help()

if __name__ == "__main__":
    main()
